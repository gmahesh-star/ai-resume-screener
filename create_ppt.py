"""
Generate a professional PowerPoint presentation for the AI Resume Screening project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour Palette ──────────────────────────────────────────────
DARK_BG       = RGBColor(0x0F, 0x0F, 0x1A)   # Deep dark navy
CARD_BG       = RGBColor(0x1A, 0x1A, 0x2E)   # Slightly lighter card
ACCENT_BLUE   = RGBColor(0x6C, 0x63, 0xFF)   # Vivid purple-blue
ACCENT_CYAN   = RGBColor(0x00, 0xD2, 0xFF)   # Bright cyan
ACCENT_GREEN  = RGBColor(0x00, 0xE6, 0x96)   # Neon green
ACCENT_PINK   = RGBColor(0xFF, 0x6B, 0x9D)   # Soft pink
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xB0, 0xB0, 0xC0)
MID_GRAY      = RGBColor(0x80, 0x80, 0x99)
ORANGE        = RGBColor(0xFF, 0xA5, 0x00)

SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ─────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, bullet_color=ACCENT_CYAN, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
        p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 – Title Slide
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Decorative elements
add_circle(slide, Inches(-1), Inches(-1), Inches(4), RGBColor(0x6C, 0x63, 0xFF))   # top-left glow
add_circle(slide, Inches(10.5), Inches(5), Inches(4), RGBColor(0x00, 0xD2, 0xFF))  # bottom-right glow

# Overlay semi-transparent card
add_rect(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.8), CARD_BG, ACCENT_BLUE, Pt(2))

# Accent bar
add_accent_line(slide, Inches(5.5), Inches(2.2), Inches(2.3), ACCENT_CYAN)

# Title
add_text_box(slide, Inches(2), Inches(2.5), Inches(9.3), Inches(1.2),
             "AI-Powered Resume Screening System", font_size=40, bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Subtitle
add_text_box(slide, Inches(2.5), Inches(3.8), Inches(8.3), Inches(1),
             "Intelligent candidate matching using NLP, Semantic Similarity & Machine Learning",
             font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom tagline
add_text_box(slide, Inches(3), Inches(5.2), Inches(7.3), Inches(0.6),
             "FastAPI  •  spaCy  •  Sentence-Transformers  •  Scikit-learn  •  SQLite",
             font_size=14, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 – Problem Statement
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.7), Inches(3), ACCENT_PINK)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(6), Inches(0.7),
             "Problem Statement", font_size=36, bold=True, color=WHITE)

problems = [
    "▸  Manual resume screening is time-consuming and error-prone",
    "▸  HR teams spend 23+ hours per hire reviewing resumes",
    "▸  Keyword-based filters miss semantically relevant candidates",
    "▸  Unconscious bias can affect manual screening processes",
    "▸  No standardized scoring system for fair comparison",
    "▸  Scaling hiring pipelines is difficult without automation",
]
add_bullet_list(slide, Inches(1), Inches(2.2), Inches(6.5), Inches(4.5),
                problems, font_size=18, color=LIGHT_GRAY)

# Right-side card
add_rect(slide, Inches(8.2), Inches(1.8), Inches(4.5), Inches(4.8), CARD_BG, ACCENT_PINK, Pt(1.5))
add_text_box(slide, Inches(8.5), Inches(2.1), Inches(4), Inches(0.5),
             "💡 The Solution", font_size=22, bold=True, color=ACCENT_PINK, alignment=PP_ALIGN.CENTER)
solution_lines = [
    "An AI-powered system that:",
    "",
    "✓ Extracts text from PDF/DOCX",
    "✓ Identifies skills automatically",
    "✓ Computes semantic similarity",
    "✓ Ranks candidates objectively",
    "✓ Provides instant results",
]
add_bullet_list(slide, Inches(8.5), Inches(2.8), Inches(4), Inches(3.5),
                solution_lines, font_size=16, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 – Project Architecture / Overview
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.7), Inches(3.5), ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7),
             "System Architecture", font_size=36, bold=True)

# Architecture boxes (flow diagram style)
components = [
    ("Frontend\n(HTML/CSS/JS)", Inches(0.8), ACCENT_CYAN),
    ("FastAPI\nBackend", Inches(3.5), ACCENT_BLUE),
    ("AI Service\n(NLP Engine)", Inches(6.2), ACCENT_GREEN),
    ("SQLite\nDatabase", Inches(8.9), ORANGE),
    ("File Parser\n(PDF/DOCX)", Inches(11.0), ACCENT_PINK),
]

box_top = Inches(2.6)
for label, left, color in components:
    add_rect(slide, left, box_top, Inches(2.1), Inches(1.3), CARD_BG, color, Pt(2))
    txBox = slide.shapes.add_textbox(left, box_top + Inches(0.2), Inches(2.1), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(label.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = color if i == 0 else LIGHT_GRAY
        p.font.bold = (i == 0)
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

# Arrows between boxes
arrow_positions = [Inches(2.9), Inches(5.6), Inches(8.3), Inches(10.85)]
for ax in arrow_positions:
    add_text_box(slide, ax, box_top + Inches(0.3), Inches(0.6), Inches(0.6),
                 "→", font_size=28, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Data flow description
add_rect(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.3), CARD_BG, RGBColor(0x30, 0x30, 0x50), Pt(1))
flow_text = [
    "① User uploads resumes (PDF/DOCX) and pastes a job description via the Web UI",
    "② FastAPI backend receives files and routes them to the File Parser for text extraction",
    "③ Extracted text is sent to the AI Service for NLP processing (spaCy + Sentence-Transformers)",
    "④ Skills are matched, TF-IDF + Semantic similarity scores are computed (70/30 weighted blend)",
    "⑤ Results are stored in SQLite, ranked by match score, and returned to the frontend",
]
add_bullet_list(slide, Inches(1.1), Inches(4.7), Inches(11.2), Inches(2),
                flow_text, font_size=14, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 – Tech Stack
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.7), Inches(2.5), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(6), Inches(0.7),
             "Technology Stack", font_size=36, bold=True)

tech_categories = [
    ("Backend", ACCENT_BLUE, [
        "Python 3.9+",
        "FastAPI (REST API framework)",
        "Uvicorn (ASGI Server)",
        "SQLAlchemy + SQLite",
    ]),
    ("AI / NLP", ACCENT_GREEN, [
        "spaCy (en_core_web_sm)",
        "Sentence-Transformers",
        "  (all-MiniLM-L6-v2)",
        "Scikit-learn (TF-IDF)",
    ]),
    ("Frontend", ACCENT_CYAN, [
        "HTML5 + CSS3 + JavaScript",
        "Glassmorphism UI Design",
        "Drag & Drop file upload",
        "Real-time progress feedback",
    ]),
    ("Utilities", ACCENT_PINK, [
        "pdfplumber (PDF parsing)",
        "python-docx (DOCX parsing)",
        "python-multipart",
        "CORS Middleware",
    ]),
]

card_width = Inches(2.8)
card_height = Inches(3.8)
start_x = Inches(0.6)
gap = Inches(0.25)

for i, (title, color, items) in enumerate(tech_categories):
    left = start_x + i * (card_width + gap)
    top = Inches(2.2)
    add_rect(slide, left, top, card_width, card_height, CARD_BG, color, Pt(2))
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), card_width - Inches(0.4), Inches(0.5),
                 title, font_size=20, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    # Divider
    add_rect(slide, left + Inches(0.3), top + Inches(0.8), card_width - Inches(0.6), Pt(2), color)
    # Items
    add_bullet_list(slide, left + Inches(0.3), top + Inches(1.1), card_width - Inches(0.5), Inches(2.5),
                    items, font_size=14, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 – AI Matching Algorithm
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.7), Inches(3), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7),
             "AI Matching Algorithm", font_size=36, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.5),
             "The system uses a dual-approach scoring model that combines traditional NLP with modern deep learning embeddings:",
             font_size=16, color=LIGHT_GRAY)

# Left card – Semantic Similarity
add_rect(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(4), CARD_BG, ACCENT_BLUE, Pt(2))
add_text_box(slide, Inches(1.1), Inches(3), Inches(5), Inches(0.5),
             "🧠 Semantic Similarity (70% weight)", font_size=20, bold=True, color=ACCENT_BLUE)
sem_items = [
    "• Uses all-MiniLM-L6-v2 Sentence Transformer",
    "• Generates 384-dim dense embeddings",
    "• Captures contextual meaning beyond keywords",
    "• Cosine similarity between resume & JD vectors",
    "• Understanding synonym relationships",
    "  (e.g., 'ML' ≈ 'Machine Learning')",
]
add_bullet_list(slide, Inches(1.1), Inches(3.7), Inches(5), Inches(3),
                sem_items, font_size=14, color=LIGHT_GRAY)

# Right card – TF-IDF
add_rect(slide, Inches(6.8), Inches(2.8), Inches(5.5), Inches(4), CARD_BG, ORANGE, Pt(2))
add_text_box(slide, Inches(7.1), Inches(3), Inches(5), Inches(0.5),
             "📊 TF-IDF Similarity (30% weight)", font_size=20, bold=True, color=ORANGE)
tfidf_items = [
    "• Traditional term frequency–inverse doc frequency",
    "• Sparse vector representation of text",
    "• Exact keyword matching strength",
    "• Cosine similarity on TF-IDF vectors",
    "• Complements semantic for precision",
    "• Catches specific technical terms",
]
add_bullet_list(slide, Inches(7.1), Inches(3.7), Inches(5), Inches(3),
                tfidf_items, font_size=14, color=LIGHT_GRAY)

# Formula bar
add_rect(slide, Inches(2.5), Inches(7.0), Inches(8.3), Inches(0.45), ACCENT_BLUE)
add_text_box(slide, Inches(2.5), Inches(7.0), Inches(8.3), Inches(0.45),
             "Final Score  =  ( Semantic × 0.70 )  +  ( TF-IDF × 0.30 )     →    0% – 100%",
             font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 – Key Features
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.7), Inches(2.5), ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(6), Inches(0.7),
             "Key Features", font_size=36, bold=True)

features = [
    ("📄", "Resume Upload", "Supports PDF & DOCX\nDrag-and-drop interface", ACCENT_CYAN),
    ("🔍", "Skill Extraction", "Auto-detects 40+ tech skills\nfrom predefined skill database", ACCENT_BLUE),
    ("🤖", "AI Matching", "Semantic + TF-IDF blend\n70/30 weighted scoring", ACCENT_GREEN),
    ("📊", "Smart Ranking", "Candidates ranked by score\nHighest match first", ORANGE),
    ("💾", "Data Persistence", "SQLite stores all results\nHistory tracking built-in", ACCENT_PINK),
    ("🌐", "Modern Web UI", "Glassmorphism design\nReal-time feedback", RGBColor(0xBB, 0x86, 0xFC)),
]

for i, (icon, title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(0.8) + col * Inches(4.1)
    top = Inches(2.2) + row * Inches(2.4)
    
    add_rect(slide, left, top, Inches(3.7), Inches(2), CARD_BG, color, Pt(1.5))
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.6), Inches(0.5),
                 icon, font_size=26, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.8), top + Inches(0.15), Inches(2.7), Inches(0.4),
                 title, font_size=18, bold=True, color=color)
    add_text_box(slide, left + Inches(0.8), top + Inches(0.7), Inches(2.7), Inches(1.1),
                 desc, font_size=13, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 – Project Structure
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.7), Inches(2.5), ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(6), Inches(0.7),
             "Project Structure", font_size=36, bold=True)

# Code-style structure
add_rect(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(5), CARD_BG, RGBColor(0x30, 0x30, 0x50), Pt(1.5))
structure_lines = [
    "Ai_resume_screening/",
    "├── backend/",
    "│   ├── main.py              # FastAPI Routes",
    "│   ├── database.py          # SQLite Database",
    "│   ├── services/",
    "│   │   └── ai_service.py    # AI/NLP Engine",
    "│   └── utils/",
    "│       └── file_parser.py   # PDF/DOCX Parser",
    "├── frontend/",
    "│   ├── index.html           # Web Interface",
    "│   ├── style.css            # Custom Styling",
    "│   └── script.js            # Client Logic",
    "├── data/",
    "│   ├── example_jd.txt       # Sample Job Description",
    "│   └── generate_samples.py  # Resume Generator",
    "├── requirements.txt",
    "└── README.md",
]
txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(5.4), Inches(4.6))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(structure_lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT_CYAN if i == 0 else LIGHT_GRAY
    p.font.name = "Consolas"
    p.font.bold = (i == 0)
    p.space_after = Pt(2)

# Right side — file descriptions
add_rect(slide, Inches(7.2), Inches(2.0), Inches(5.3), Inches(5), CARD_BG, ACCENT_BLUE, Pt(1.5))
add_text_box(slide, Inches(7.5), Inches(2.2), Inches(4.8), Inches(0.4),
             "Module Descriptions", font_size=20, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

module_descs = [
    ("main.py", "REST API endpoints: /upload-resume, /results, /results (DELETE). Handles file uploads, CORS, and startup events."),
    ("ai_service.py", "Core NLP engine: text preprocessing with spaCy, skill extraction from SKILLS_DB, TF-IDF & Sentence-Transformer similarity scoring."),
    ("database.py", "SQLite persistence layer: stores filenames, scores, matched/missing skills. Supports history retrieval and clearing."),
    ("file_parser.py", "Extracts raw text from uploaded PDF (pdfplumber) and DOCX (python-docx) files."),
]

y_offset = Inches(2.8)
for name, desc in module_descs:
    add_text_box(slide, Inches(7.5), y_offset, Inches(4.8), Inches(0.3),
                 name, font_size=14, bold=True, color=ACCENT_CYAN, font_name="Consolas")
    add_text_box(slide, Inches(7.5), y_offset + Inches(0.3), Inches(4.8), Inches(0.7),
                 desc, font_size=12, color=LIGHT_GRAY)
    y_offset += Inches(1.05)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 – Frontend UI
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.7), Inches(2.5), ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(6), Inches(0.7),
             "Frontend Design", font_size=36, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.5),
             "A sleek, modern glassmorphism interface built with vanilla HTML/CSS/JS — no external frameworks required.",
             font_size=16, color=LIGHT_GRAY)

# UI Feature cards
ui_features = [
    ("Navigation", "Sticky nav bar with smooth\nscroll between Upload,\nResults, and History sections", ACCENT_CYAN),
    ("Upload Zone", "Drag-and-drop area with\nreal-time file list preview,\nchar counter for JD textarea", ACCENT_BLUE),
    ("Results View", "Animated score cards with\ncolor-coded match bars,\nskill chips (green/red)", ACCENT_GREEN),
    ("Loading State", "Multi-step progress indicator\nshowing: Extract → Parse →\nMatch → Rank stages", ACCENT_PINK),
]

for i, (title, desc, color) in enumerate(ui_features):
    left = Inches(0.8) + i * Inches(3.15)
    top = Inches(2.8)
    add_rect(slide, left, top, Inches(2.9), Inches(2.5), CARD_BG, color, Pt(2))
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(2.5), Inches(0.4),
                 title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_rect(slide, left + Inches(0.3), top + Inches(0.7), Inches(2.3), Pt(2), color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.9), Inches(2.5), Inches(1.4),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Design highlights
add_rect(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.5), CARD_BG, RGBColor(0x30, 0x30, 0x50), Pt(1))
design_items = [
    "🎨 Dark theme with glassmorphism cards  •  Animated background glows & grid pattern  •  Inter + JetBrains Mono typography",
    "⚡ Animated spinner with 4-step progress  •  Hover effects on all interactive elements  •  Responsive layout for all screen sizes",
    "🏆 Color-coded match scores: Green (>70%), Yellow (40-70%), Red (<40%)  •  Skill chips with matched/missing indicators",
]
add_bullet_list(slide, Inches(1.1), Inches(5.8), Inches(11.2), Inches(1.2),
                design_items, font_size=13, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 – Example Output / Demo
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.7), Inches(3), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7),
             "Example Output — Demo Results", font_size=36, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.5),
             'Job Description: "Senior AI Engineer — Python, ML, NLP, FastAPI, Docker..."',
             font_size=14, color=LIGHT_GRAY)

# Candidate cards
candidates = [
    ("🥇", "Rank #1", "Jane_Doe_Resume.docx", "~85%", ACCENT_GREEN,
     "python, machine learning, nlp, fastapi, spacy, git, docker", "sql"),
    ("🥈", "Rank #2", "John_Smith_Resume.docx", "~35%", ORANGE,
     "python", "machine learning, nlp, docker, fastapi, spacy, git"),
    ("🥉", "Rank #3", "Alice_Jones_Resume.docx", "~15%", ACCENT_PINK,
     "—", "Most key indicators missing"),
]

for i, (medal, rank, filename, score, color, matched, missing) in enumerate(candidates):
    left = Inches(0.8) + i * Inches(4.15)
    top = Inches(2.6)
    add_rect(slide, left, top, Inches(3.8), Inches(4.5), CARD_BG, color, Pt(2))
    
    # Medal + Rank
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.7), Inches(0.5),
                 medal, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.9), top + Inches(0.15), Inches(2.5), Inches(0.4),
                 rank, font_size=20, bold=True, color=color)
    
    # Filename
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7), Inches(3.2), Inches(0.3),
                 filename, font_size=13, color=LIGHT_GRAY, font_name="Consolas")
    
    # Score
    add_rect(slide, left + Inches(0.3), top + Inches(1.2), Inches(3.2), Inches(0.6), RGBColor(0x12, 0x12, 0x20), color, Pt(1))
    add_text_box(slide, left + Inches(0.3), top + Inches(1.2), Inches(3.2), Inches(0.6),
                 f"Match Score: {score}", font_size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    
    # Matched skills
    add_text_box(slide, left + Inches(0.3), top + Inches(2.1), Inches(3.2), Inches(0.3),
                 "✅ Matched Skills:", font_size=12, bold=True, color=ACCENT_GREEN)
    add_text_box(slide, left + Inches(0.3), top + Inches(2.5), Inches(3.2), Inches(0.7),
                 matched, font_size=11, color=LIGHT_GRAY)
    
    # Missing skills
    add_text_box(slide, left + Inches(0.3), top + Inches(3.3), Inches(3.2), Inches(0.3),
                 "❌ Missing Skills:", font_size=12, bold=True, color=ACCENT_PINK)
    add_text_box(slide, left + Inches(0.3), top + Inches(3.7), Inches(3.2), Inches(0.7),
                 missing, font_size=11, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 – API Endpoints
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.7), Inches(2.5), ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(6), Inches(0.7),
             "API Endpoints", font_size=36, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.5),
             "RESTful API built with FastAPI — auto-generated Swagger docs at /docs",
             font_size=16, color=LIGHT_GRAY)

endpoints = [
    ("POST", "/upload-resume", "Upload resumes + job description for AI analysis",
     "Body: job_description (form), files[] (multipart)\nResponse: Ranked list of candidates with scores & skills", ACCENT_GREEN),
    ("GET", "/results", "Retrieve all past screening results",
     "Response: All stored results ordered by match_score DESC\nIncludes: filename, score, matched/missing skills, rank", ACCENT_CYAN),
    ("DELETE", "/results", "Clear all stored screening history",
     "Response: { status: 'success', message: 'Results cleared' }\nResets the SQLite database table", ACCENT_PINK),
]

for i, (method, path, desc, details, color) in enumerate(endpoints):
    top = Inches(2.6) + i * Inches(1.6)
    add_rect(slide, Inches(0.8), top, Inches(11.7), Inches(1.4), CARD_BG, color, Pt(1.5))
    
    # Method badge
    add_rect(slide, Inches(1.0), top + Inches(0.2), Inches(1.2), Inches(0.4), color)
    add_text_box(slide, Inches(1.0), top + Inches(0.2), Inches(1.2), Inches(0.4),
                 method, font_size=14, bold=True, color=DARK_BG, alignment=PP_ALIGN.CENTER)
    
    # Path
    add_text_box(slide, Inches(2.4), top + Inches(0.15), Inches(2.5), Inches(0.4),
                 path, font_size=16, bold=True, color=color, font_name="Consolas")
    
    # Description
    add_text_box(slide, Inches(5.0), top + Inches(0.15), Inches(7.2), Inches(0.4),
                 desc, font_size=14, color=WHITE)
    
    # Details
    add_text_box(slide, Inches(2.4), top + Inches(0.65), Inches(9.8), Inches(0.7),
                 details, font_size=11, color=MID_GRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 – Setup & Installation
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.7), Inches(3), ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(6), Inches(0.7),
             "Setup & Installation", font_size=36, bold=True)

steps = [
    ("1", "Create Virtual Environment", "python -m venv venv\nvenv\\Scripts\\activate", ACCENT_CYAN),
    ("2", "Install Dependencies", "pip install -r requirements.txt", ACCENT_BLUE),
    ("3", "Download NLP Model", "python -m spacy download en_core_web_sm", ACCENT_GREEN),
    ("4", "Start Backend Server", "uvicorn backend.main:app --reload", ORANGE),
    ("5", "Open Frontend", "Open frontend/index.html in browser\nor: python -m http.server 8080", ACCENT_PINK),
]

for i, (num, title, cmd, color) in enumerate(steps):
    top = Inches(2.0) + i * Inches(1.05)
    
    # Step number circle
    add_circle(slide, Inches(0.9), top + Inches(0.05), Inches(0.5), color)
    add_text_box(slide, Inches(0.9), top + Inches(0.05), Inches(0.5), Inches(0.5),
                 num, font_size=18, bold=True, color=DARK_BG, alignment=PP_ALIGN.CENTER)
    
    # Title
    add_text_box(slide, Inches(1.6), top, Inches(3), Inches(0.4),
                 title, font_size=18, bold=True, color=color)
    
    # Command box  
    add_rect(slide, Inches(5), top, Inches(7.5), Inches(0.85), CARD_BG, RGBColor(0x30, 0x30, 0x50), Pt(1))
    add_text_box(slide, Inches(5.2), top + Inches(0.1), Inches(7.1), Inches(0.7),
                 cmd, font_size=13, color=ACCENT_CYAN, font_name="Consolas")

# ══════════════════════════════════════════════════════════════════
# SLIDE 12 – Future Enhancements
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.7), Inches(3.5), RGBColor(0xBB, 0x86, 0xFC))
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7),
             "Future Enhancements", font_size=36, bold=True)

enhancements = [
    ("🧪", "Resume Improvement Suggestions", "Provide AI-powered feedback to candidates on how to improve their resume for a given role.", ACCENT_CYAN),
    ("🎓", "Education & Experience Detection", "Extract and score educational qualifications, certifications, and years of experience.", ACCENT_GREEN),
    ("📈", "Advanced Analytics Dashboard", "Visual charts showing skill distribution, score trends, and batch comparison statistics.", ACCENT_BLUE),
    ("🔐", "User Authentication", "Secure login, role-based access control, and per-recruiter history tracking.", ACCENT_PINK),
    ("☁️", "Cloud Deployment", "Docker containerization, CI/CD pipeline, and deployment to AWS/GCP/Azure.", ORANGE),
    ("🤖", "LLM Integration", "Use GPT/Gemini for generating interview questions tailored to each candidate's profile.", RGBColor(0xBB, 0x86, 0xFC)),
]

for i, (icon, title, desc, color) in enumerate(enhancements):
    col = i % 2
    row = i // 2
    left = Inches(0.8) + col * Inches(6.3)
    top = Inches(2.2) + row * Inches(1.65)
    
    add_rect(slide, left, top, Inches(5.9), Inches(1.4), CARD_BG, color, Pt(1.5))
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), Inches(0.6), Inches(0.5),
                 icon, font_size=24, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.8), top + Inches(0.1), Inches(4.8), Inches(0.35),
                 title, font_size=16, bold=True, color=color)
    add_text_box(slide, left + Inches(0.8), top + Inches(0.55), Inches(4.8), Inches(0.75),
                 desc, font_size=12, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 13 – Thank You
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_circle(slide, Inches(-1.5), Inches(-1.5), Inches(5), ACCENT_BLUE)
add_circle(slide, Inches(10), Inches(4.5), Inches(5), ACCENT_CYAN)

add_rect(slide, Inches(2.5), Inches(1.8), Inches(8.3), Inches(4.2), CARD_BG, ACCENT_BLUE, Pt(2))

add_accent_line(slide, Inches(5.5), Inches(2.5), Inches(2.3), ACCENT_CYAN)

add_text_box(slide, Inches(3), Inches(2.8), Inches(7.3), Inches(1),
             "Thank You!", font_size=48, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(3.9), Inches(7.3), Inches(0.7),
             "AI-Powered Resume Screening System",
             font_size=22, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(4.6), Inches(7.3), Inches(0.7),
             "Built with Python • FastAPI • spaCy • Sentence-Transformers",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(5.3), Inches(7.3), Inches(0.5),
             "Questions?",
             font_size=20, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────
output_path = r"d:\Ai_resume_screening\AI_Resume_Screening_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
